import logging
import PyPDF2
import pandas as pd
import docx
from pptx import Presentation
from io import BytesIO
import os
import tempfile
import shutil
import re
import aiohttp
from bs4 import BeautifulSoup
from typing import Tuple, Optional, List, Dict, Any
from message_processing import process_and_send_ai_response

# URL regular expression for detecting links in text
URL_REGEX = r'https?://[^\s<>"]+|www\.[^\s<>"]+'


async def process_content_after_timeout(user_id, message, state):
    """Process collected content (documents and URLs) after a timeout"""
    import asyncio
    from config import message_timers, collected_content
    
    # Wait for 1 second
    await asyncio.sleep(1)
    
    # Clear timer reference
    message_timers[user_id] = None
    
    # Check if we have content to process
    if user_id in collected_content and collected_content[user_id]:
        content_count = len(collected_content[user_id])
        
        # No user feedback will be sent here, it's handled elsewhere
        
        # Process the content now that we've gathered all items
        await analyze_content(message, state)

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            
            # Limit to first 50 pages for very large PDFs
            if num_pages > 50:
                logging.info(f"PDF has {num_pages} pages, limiting extraction to first 50 pages")
                num_pages = 50
                
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                if page_text:
                    text += f"--- Page {page_num + 1} ---\n{page_text}\n\n"
                    
            if not text.strip():
                text = "[PDF does not contain extractable text, possibly contains scanned images]"
                
    except Exception as e:
        logging.error(f"Error in PDF extraction: {e}")
        text = f"[Error extracting text from PDF: {str(e)}]"
    
    # Limit text length for very large documents
    if len(text) > 100000:  # ~25 pages of text
        text = text[:100000] + "\n\n[Text truncated due to large document size...]"
    
    return text


def extract_text_from_docx(file_path):
    """Extract text from a DOCX file"""
    text = ""
    try:
        doc = docx.Document(file_path)
        
        # Get document paragraphs
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        # Get text from tables
        tables_text = []
        for table in doc.tables:
            table_text = ""
            for i, row in enumerate(table.rows):
                row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                if row_text:
                    table_text += row_text + "\n"
            if table_text:
                tables_text.append(table_text)
        
        # Combine all text
        if paragraphs:
            text += "\n".join(paragraphs) + "\n\n"
        if tables_text:
            text += "--- Таблицы ---\n" + "\n".join(tables_text)
        
        if not text.strip():
            text = "[Документ не содержит текста]"
            
    except Exception as e:
        logging.error(f"Error in DOCX extraction: {e}")
        text = f"[Ошибка извлечения текста из DOCX: {str(e)}]"
    
    return text


def extract_text_from_xlsx(file_path):
    """Extract text from an XLSX file"""
    text = ""
    try:
        # Read Excel file
        excel = pd.ExcelFile(file_path)
        sheet_names = excel.sheet_names
        
        # Limit to first 5 sheets for large files
        if len(sheet_names) > 5:
            logging.info(f"Excel has {len(sheet_names)} sheets, limiting extraction to first 5")
            sheet_names = sheet_names[:5]
        
        for sheet_name in sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Skip empty sheets
            if df.empty:
                continue
                
            # Limit rows for very large sheets (extract first 100 rows)
            if len(df) > 100:
                logging.info(f"Sheet {sheet_name} has {len(df)} rows, limiting extraction to first 100")
                df = df.head(100)
            
            # Convert dataframe to string
            sheet_text = f"--- Лист: {sheet_name} ---\n"
            
            # Try to get a reasonable string representation
            try:
                sheet_text += df.to_string(index=False) + "\n\n"
            except Exception:
                try:
                    # Alternative method with fewer columns
                    if len(df.columns) > 10:
                        sheet_text += df.iloc[:, :10].to_string(index=False)
                        sheet_text += "\n[Показаны только первые 10 колонок]\n\n"
                    else:
                        sheet_text += str(df) + "\n\n"
                except Exception as e:
                    sheet_text += f"[Не удалось извлечь данные из листа: {str(e)}]\n\n"
            
            text += sheet_text
        
        if not text.strip():
            text = "[Excel файл не содержит данных или не может быть прочитан]"
            
    except Exception as e:
        logging.error(f"Error in XLSX extraction: {e}")
        text = f"[Ошибка извлечения данных из Excel: {str(e)}]"
    
    return text


def extract_text_from_pptx(file_path):
    """Extract text from a PPTX file with enhanced error handling"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import zipfile
    import xml.etree.ElementTree as ET
    
    def safe_extract_text(shape):
        """Safely extract text from a shape with multiple fallback methods"""
        try:
            # Try standard text extraction first
            if hasattr(shape, 'text') and shape.text and shape.text.strip():
                return shape.text.strip() + "\n"
                
            # Try text frame extraction
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = ''
                for paragraph in shape.text_frame.paragraphs:
                    if hasattr(paragraph, 'text') and paragraph.text.strip():
                        text += paragraph.text.strip() + "\n"
                    # Try runs if direct text is empty
                    elif hasattr(paragraph, 'runs'):
                        for run in paragraph.runs:
                            if hasattr(run, 'text') and run.text.strip():
                                text += run.text.strip() + "\n"
                if text:
                    return text
                    
            # Handle tables
            if hasattr(shape, 'has_table') and shape.has_table:
                table_text = ''
                try:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if hasattr(cell, 'text_frame'):
                                cell_text = ' '.join([p.text for p in cell.text_frame.paragraphs 
                                                   if hasattr(p, 'text') and p.text.strip()])
                                if cell_text.strip():
                                    row_text.append(cell_text.strip())
                        if row_text:
                            table_text += ' | '.join(row_text) + "\n"
                    if table_text:
                        return table_text
                except Exception as e:
                    logging.warning(f"Error extracting table: {e}")
            
            # Handle group shapes
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = ''
                if hasattr(shape, 'shapes'):
                    for sub_shape in shape.shapes:
                        group_text += safe_extract_text(sub_shape) or ''
                return group_text
                
        except Exception as e:
            logging.warning(f"Error in safe_extract_text: {e}")
        
        return ''
    
    # First, try to validate the PPTX file
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            # Check if this is a valid zip file
            if 'ppt/presentation.xml' not in z.namelist():
                return "[Неверный формат файла презентации]"
    except Exception as e:
        return f"[Ошибка чтения файла: {str(e)}]"
    
    # Now try to extract text
    text = ""
    try:
        prs = Presentation(file_path)
        num_slides = len(prs.slides)
        
        # Limit to first 30 slides for very large presentations
        max_slides = min(30, num_slides)
        if num_slides > 30:
            logging.info(f"Presentation has {num_slides} slides, limiting extraction to first 30")
        
        for i, slide in enumerate(prs.slides[:max_slides]):
            try:
                slide_text = f"--- Слайд {i+1} ---\n"
                shape_texts = []
                
                # Process shapes with error handling for each
                for shape in slide.shapes:
                    try:
                        shape_text = safe_extract_text(shape)
                        if shape_text:
                            shape_texts.append(shape_text)
                    except Exception as shape_error:
                        logging.warning(f"Error in shape {i+1}-{slide.shapes.index(shape)}: {shape_error}")
                
                # Add all successfully extracted shape texts
                if shape_texts:
                    slide_text += '\n'.join(shape_texts) + '\n'
                
                # Try to add notes if available
                try:
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_text += "\nПримечания:\n" + notes_text + "\n"
                except Exception as notes_error:
                    logging.warning(f"Error reading notes for slide {i+1}: {notes_error}")
                
                # Only add non-empty slides
                if slide_text.strip() != f"--- Слайд {i+1} ---\n":
                    text += slide_text + "\n"
            
            except Exception as slide_error:
                logging.error(f"Error processing slide {i+1}: {slide_error}")
                continue
                
        if not text.strip():
            text = "[Не удалось извлечь текст из презентации]"
            
    except Exception as e:
        logging.error(f"Error in PPTX extraction: {e}")
        text = f"[Ошибка извлечения текста из презентации: {str(e)}]"
    
    return text


async def collect_content(message, state, content_info):
    """Collect content (documents, URLs) with a timer for batch processing"""
    import asyncio
    from config import message_timers, collected_content
    
    # Get user ID
    user_id = message.from_user.id
    
    # Ensure the user has an entry in collected_content
    if user_id not in collected_content:
        collected_content[user_id] = []
    
    # Добавляем контент в коллекцию (не очищаем предыдущий контент)
    collected_content[user_id].append(content_info)
    
    # If there's already a timer, cancel it and create a new one to process the entire content set
    if user_id in message_timers and message_timers[user_id] is not None:
        try:
            message_timers[user_id].cancel()
        except Exception as e:
            logging.error(f"Ошибка отмены таймера для пользователя {user_id}: {e}")
    
    # Создаем новый таймер (1 секунда) - в течение этого времени собираем документы и ссылки
    timer_task = asyncio.create_task(process_content_after_timeout(user_id, message, state))
    message_timers[user_id] = timer_task
    
    # Не отправляем уведомление о получении документа - это будет сделано в process_content_after_timeout после таймаута
    # This is needed to process a batch of documents and links forwarded by the user within a short time period
    pass


async def extract_url_from_text(text: str) -> Optional[str]:
    """
    Extracts URL from message text
    
    Args:
        text: Message text
        
    Returns:
        URL or None if URL is not found
    """
    if not text:
        return None
        
    # Search for URL in text
    match = re.search(URL_REGEX, text)
    if match:
        url = match.group(0)
        # If URL starts with www., add http:// protocol
        if url.startswith('www.'):
            url = 'http://' + url
        return url
    
    return None

async def parse_url_content(url: str, max_length: int = 15000) -> Tuple[str, str, Optional[str]]:
    """
    Parses URL content and returns title, text and image
    
    Args:
        url: URL to parse
        max_length: Maximum length of extracted text
        
    Returns:
        Tuple (title, text, image URL)
    """
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=10) as response:
                if response.status != 200:
                    return url, f"Failed to load page content. Status code: {response.status}", None
                
                html = await response.text()
                
                # Parse HTML with BeautifulSoup
                soup = BeautifulSoup(html, 'html.parser')
                
                # Remove unnecessary elements
                for element in soup.find_all(['script', 'style', 'iframe', 'nav', 'footer']):
                    element.decompose()
                
                # Extract title
                title = soup.title.text.strip() if soup.title else "No title"
                
                # Extract main text
                paragraphs = []
                for p in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li']):
                    text = p.get_text().strip()
                    if text and len(text) > 10:  # Ignore short paragraphs
                        paragraphs.append(text)
                
                # If text is scarce, look for main text in other tags
                if len(paragraphs) < 3:
                    for div in soup.find_all(['div', 'article', 'section', 'main']):
                        text = div.get_text().strip()
                        if text and len(text) > 50:  # Only text with decent length
                            paragraphs.append(text)
                
                # Combine text, limit length
                text = "\n\n".join(paragraphs)
                if len(text) > max_length:
                    text = text[:max_length] + "... [text truncated]"
                
                # If text is too short or empty, return error message
                if len(text) < 50:
                    return title, f"Failed to extract enough text from site {url}", None
                
                # Find the first large image (for preview)
                img_url = None
                for img in soup.find_all('img'):
                    src = img.get('src')
                    if src and not src.startswith('data:'):
                        # Convert relative URL to absolute
                        if not src.startswith(('http://', 'https://')):
                            if src.startswith('/'):
                                # Relative from root
                                base_url = '/'.join(url.split('/')[:3])  # http://domain.com
                                src = base_url + src
                            else:
                                # Relative from current path
                                src = url.rsplit('/', 1)[0] + '/' + src
                        img_url = src
                        break
                
                return title, text, img_url
    
    except Exception as e:
        logging.error(f"Error parsing URL {url}: {str(e)}")
        return url, f"Error parsing URL: {str(e)}", None

async def handle_document(message, state, bot):
    """Handle document uploads with timer-based batching"""
    import os
    from config import collected_content
    
    # Handle both document and photo types
    if message.document:
        # Extract document information
        doc = message.document
        file_id = doc.file_id
        file_name = doc.file_name
        file_size = doc.file_size
        mime_type = doc.mime_type
        
        logging.info(f"Received document: {file_name}, mime_type: {mime_type}, size: {file_size} bytes")
        
        # Check supported file types
        supported_types = {
            'application/pdf': {'ext': 'pdf', 'processor': extract_text_from_pdf},
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': {'ext': 'docx', 'processor': extract_text_from_docx},
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': {'ext': 'xlsx', 'processor': extract_text_from_xlsx},
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': {'ext': 'pptx', 'processor': extract_text_from_pptx}
        }
        
        # Also check by extension for cases where mime_type is generic
        if file_name:
            ext = file_name.split('.')[-1].lower() if '.' in file_name else ''
            extension_map = {
                'pdf': {'mime': 'application/pdf', 'processor': extract_text_from_pdf},
                'docx': {'mime': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'processor': extract_text_from_docx},
                'xlsx': {'mime': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'processor': extract_text_from_xlsx},
                'pptx': {'mime': 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'processor': extract_text_from_pptx}
            }
            if ext in extension_map and mime_type not in supported_types:
                # Use extension mapping if mime type not recognized
                mime_type = extension_map[ext]['mime']
        
        if mime_type not in supported_types:
            await message.reply(f"This file type is not supported. Please use PDF, DOCX, XLSX, or PPTX files.")
            return
        
        # Download the file
        file = await bot.get_file(file_id)
        file_path = file.file_path
        
        # Create directory if not exists
        os.makedirs('/tmp/telegram_doc_analyzer', exist_ok=True)
        local_path = f"/tmp/telegram_doc_analyzer/{file_name}"
        
        # Download the file
        await bot.download_file(file_path, local_path)
        
        # Get processor function
        processor = supported_types[mime_type]['processor']
        
        # Create content info
        content_info = {
            'type': 'document',
            'file_id': file_id,
            'file_name': file_name,
            'mime_type': mime_type,
            'file_path': local_path,
            'processor': processor.__name__
        }
        
        # Also check for URLs in caption
        if message.caption:
            url = await extract_url_from_text(message.caption)
            if url:
                # Note: We'll process URL later together with document
                content_info['url'] = url
                logging.info(f"Found URL in document caption: {url}")
        
        # Collect for batched processing
        await collect_content(message, state, content_info)


async def handle_url(message, state):
    """
    Handle URL messages with timer-based batching
    
    Args:
        message: The message containing the URL
        state: FSM state
    """
    # Process URL from message
    url = await extract_url_from_text(message.text or message.caption or "")
    
    if not url:
        return
    
    logging.info(f"Handling URL: {url}")
    
    # Create content info for URL
    content_info = {
        'type': 'url',
        'url': url,
        'message_text': message.text or message.caption or ""
    }
    
    # Collect for batched processing
    await collect_content(message, state, content_info)
    return True  # URL was found and collected

async def analyze_content(message, state):
    """Analyze a batch of collected content (documents and URLs)"""
    # Import necessary modules at the beginning to avoid import errors
    import traceback
    from config import collected_content
    from claude_utils import get_claude_completion, active_sessions, close_session
    from config import processing_users
    from redis_storage import RedisStorage
    
    user_id = message.from_user.id
    
    # Check if we have content and not already processing for this user
    if user_id not in collected_content or not collected_content[user_id] or user_id in processing_users:
        return
    
    # Add user to processing list to prevent concurrent requests
    processing_users.add(user_id)
    
    # According to new requirements, forwarded messages, documents, and URLs should ALWAYS start a new chat
    from claude_utils import close_session, create_session
    from redis_storage import RedisStorage
    redis_storage = RedisStorage()
    user_rules = await redis_storage.get_user_rules(user_id)
    
    # Document/URL analysis is treated as a new conversation context
    # This allows users to engage with the document contents in a focused conversation
    
    # Explicitly close any existing session and create a new one for content analysis
    logging.info(f"Documents/URLs detected - starting new chat for user {user_id}")
    
    # Always close the existing session first
    if user_id in active_sessions:
        logging.info(f"Closing existing session for user {user_id} before document analysis")
        await close_session(user_id)
    
    # Create a new session - will be done automatically in the content analysis flow
    
    # Clear state data but preserve rules
    await state.clear()
    if user_rules:
        await state.update_data(rules=user_rules)
    
    # Get state data
    state_data = await state.get_data()
    
    # Get the content types
    content_types = set(item.get('type', '') for item in collected_content[user_id])
    
    # Inform user that we're analyzing content
    if 'url' in content_types and 'document' in content_types:
        working_msg = await message.answer("Анализирую документы и ссылки...")
    elif 'url' in content_types:
        working_msg = await message.answer("Анализирую ссылки...")
    else:
        working_msg = await message.answer("Анализирую документы...")
    
    # Extract text from each content item
    all_content_text = ""
    for content_info in collected_content[user_id]:
        content_type = content_info.get('type', '')
        if content_type == 'document':
            try:
                # Get the processor function name from stored info
                processor_name = content_info['processor']
                
                # Get actual function reference
                processor_functions = {
                    'extract_text_from_pdf': extract_text_from_pdf,
                    'extract_text_from_docx': extract_text_from_docx,
                    'extract_text_from_xlsx': extract_text_from_xlsx,
                    'extract_text_from_pptx': extract_text_from_pptx
                }
                
                if processor_name not in processor_functions:
                    raise ValueError(f"Unknown processor function: {processor_name}")
                
                processor = processor_functions[processor_name]
                
                # Extract text from the document
                file_path = content_info['file_path']
                file_name = content_info['file_name']
                extracted_text = processor(file_path)
                
                # Add to combined text
                if all_content_text:
                    all_content_text += "\n\n" + "-" * 40 + "\n\n"
                all_content_text += f"=== Document: {file_name} ===\n\n{extracted_text}\n\n"
                
                # Check if there's an URL in the document caption
                if 'url' in content_info:
                    url = content_info['url']
                    title, url_text, _ = await parse_url_content(url)
                    if all_content_text:
                        all_content_text += "\n\n" + "-" * 40 + "\n\n"
                    all_content_text += f"=== URL in document caption: {title} ===\n\n{url_text}\n\n"
                
            except Exception as e:
                logging.error(f"Error processing document {content_info.get('file_name', 'unknown')}: {str(e)}")
                if all_content_text:
                    all_content_text += "\n\n" + "-" * 40 + "\n\n"
                all_content_text += f"=== Error processing document {content_info.get('file_name', 'unknown')} ===\n\n{str(e)}\n\n"
        
        elif content_type == 'url':
            try:
                url = content_info['url']
                title, url_text, _ = await parse_url_content(url)
                
                # Add to combined text
                if all_content_text:
                    all_content_text += "\n\n" + "-" * 40 + "\n\n"
                all_content_text += f"=== URL: {title} ===\n\n{url_text}\n\n"
            except Exception as e:
                logging.error(f"Error processing URL {content_info.get('url', 'unknown')}: {str(e)}")
                if all_content_text:
                    all_content_text += "\n\n" + "-" * 40 + "\n\n"
                all_content_text += f"=== Error processing URL {content_info.get('url', 'unknown')} ===\n\n{str(e)}\n\n"
    
    # Get forwarded text messages from state and original message text
    data = await state.get_data()
    forwarded_text_messages = data.get('forwarded_messages', [])
    original_message_text = message.text or message.caption or ""
    
    # Determine if there is any real text content (not just a URL)
    # Check if the original message text is just a URL
    is_just_url = False
    original_text_stripped = original_message_text.strip()
    if original_text_stripped:
        url_match = re.search(URL_REGEX, original_text_stripped)
        if url_match and len(url_match.group(0)) == len(original_text_stripped):
            is_just_url = True
            logging.info("Original text is just a URL, not considering it as accompanying text")
    
    # Check for forwarded text messages or meaningful text in the message
    has_forwarded_texts = bool(forwarded_text_messages)
    has_text_content = len(original_text_stripped) > 5 and not is_just_url
    
    # Mixed mode = document/link analysis + variants ONLY when there's text content
    # If there's only documents/links without text - just give a summary (no variants)
    if has_forwarded_texts or has_text_content:
        # If there's meaningful text, use mixed mode with variants
        content_analysis_mode = True
        variant_mode = True 
        mixed_content_mode = True
        logging.info(f"Text content detected with documents/links. Setting mixed_content_mode=True with variants")
    else:
        # If there's ONLY documents/links (no text), only give summary, no variants
        content_analysis_mode = True
        variant_mode = False
        mixed_content_mode = False
        logging.info(f"Only documents/links without text. Using content analysis mode without variants")
    
    # Save modes in user state for subsequent messages
    await state.update_data(
        variant_mode=variant_mode, 
        content_analysis_mode=content_analysis_mode,
        mixed_content_mode=mixed_content_mode,
        document_mode=True,
        is_continuation=True
    )
    
    # Save document texts in state for context of subsequent messages
    await state.update_data(document_texts=all_content_text[:10000])
    
    # Create system message based on text presence
    # In any case, first analyze documents and provide factual information
    # Create base system message depending on whether there are text messages
    if has_forwarded_texts or has_text_content:
        # If there's text or a question, prepare system prompt for mixed mode
        system_message = (
            "Вы - ассистент, который анализирует документы и предоставляет помощь пользователю. "
            "КРАЙНЕ ВАЖНО: Все ответы ВСЕГДА ДОЛЖНЫ БЫТЬ НА РУССКОМ ЯЗЫКЕ, независимо от языка оригинальных документов. "
            "Если документы на английском или любом другом языке, переведите ключевую информацию на русский язык. "
            "Ваша задача состоит из двух частей: "
            "1) Сначала предоставьте ФАКТИЧЕСКОЕ краткое резюме содержимого документов и ссылок. "
            "2) Затем, на основе проанализированных документов, предложите несколько вариантов ответов или действий, "
            "которые соответствуют контексту пересланных сообщений или запросу пользователя. "
            "Между резюме и вариантами ответов используйте маркер ### "
            "Дайте 2-4 различных варианта ответа, каждый начните с ### "
        )
        
        # Add user rules if available
        if user_rules:
            system_message += f"\n\nUser rules: {user_rules}"
            
        # Create message based on whether there are forwarded text messages
        user_message = "Analyze the following documents and links, then suggest response options:\n\n"
        
        # Add user message text or forwarded messages
        if has_text_content and original_message_text.strip():
            user_message = f"User message: {original_message_text}\n\n" + user_message
            
        # Add forwarded text messages if they exist
        if has_forwarded_texts:
            forwarded_text = "\n\n=== Forwarded messages: ===\n"
            for msg in forwarded_text_messages:
                if 'text' in msg and msg['text']:
                    forwarded_text += f"\n- {msg.get('sender', 'User')}: {msg['text']}"
            user_message += forwarded_text + "\n\n"
            
        # Add document and link content
        user_message += all_content_text
        
        messages = [
            {"role": "system", "content": system_message},
            {"role": "user", "content": user_message}
        ]
    else:
        # Standard mode for documents without text - only concise summary
        system_message = (
            "Вы - ассистент, который объективно анализирует документы. "
            "Ваша задача - предоставить СТРОГО ФАКТИЧЕСКОЕ краткое резюме (3-5 предложений) того, что СОДЕРЖИТСЯ в документе. "
            "НИКОГДА не добавляйте собственные интерпретации, предположения или информацию, которой нет в документе. "
            "НИКОГДА не пишите от первого лица, не используйте фразы как 'билеты закупил', 'я купил', и подобные. "
            "КРАЙНЕ ВАЖНО: Резюме ВСЕГДА ДОЛЖНО БЫТЬ НА РУССКОМ ЯЗЫКЕ, независимо от исходного языка документа. "
            "Если документы на английском или любом другом языке, переведите ключевые моменты на русский для резюме. "
            "Сосредоточьтесь ТОЛЬКО на фактической информации из документов. "
            "Если актуально, включите фактические детали: КОГДА, ГДЕ, В КАКОЕ ВРЕМЯ, КТО (без указания конкретных имен)."
        )
        
        # Add user rules if available
        if user_rules:
            system_message += f"\n\nПравила пользователя: {user_rules}"
            
        messages = [
            {"role": "system", "content": system_message},
            {"role": "user", "content": f"Проанализируйте следующие документы и предоставьте краткое резюме:\n\n{all_content_text}"}
        ]
    
    # Get AI completion with persistent user session
    response = await get_claude_completion(messages, max_tokens=2000, message=message, user_id=user_id)

    # Delete working message if exists (AFTER getting Claude response)
    if working_msg:
        await working_msg.delete()

    # Add user query, document text, and AI response to conversation history
    if response:
        # Используем функции истории разговоров из claude_utils
        
        # Log key information for debugging
        logging.info(f"Document analysis completed. Saving to conversation history. response_length={len(response)}")
        
        # Save user query if present
        original_message_text = message.text or message.caption or ""
        if original_message_text.strip():
            # История разговоров теперь управляется Claude через active_sessions
            logging.info(f"Added user text to conversation: {original_message_text[:50]}...")
        
        # Save document as separate user message - critically important for subsequent questions
        document_message = f"[Документ]\n{all_content_text[:2000]}"
        if len(all_content_text) > 2000:
            document_message += "\n[документ слишком большой, показан фрагмент]"
        
        # История документов автоматически сохраняется в active_sessions Claude
        logging.info(f"Added document content to conversation history, length={len(document_message)}")
        
        # Save AI response to conversation history
        # Ответ AI автоматически сохраняется в active_sessions Claude
        logging.info("Added AI response to conversation history")
        
        # Проверка наличия текстовых сообщений среди пересланных для определения режима вариантов ответов
        # Получим данные из state, чтобы проверить, были ли пересланы текстовые сообщения
        data = await state.get_data()
        forwarded_text_messages = data.get('forwarded_messages', [])
        has_text_messages = bool(forwarded_text_messages)
        
        # Если есть пересланные текстовые сообщения или документы с текстом, включаем режим вариантов
        # Mixed mode: content_analysis_mode=True (анализ документов) + variant_mode=True (варианты ответов)
        if has_text_messages:
            logging.info("Detected mixed mode: content + text messages, enabling variants mode")
            await state.update_data(variant_mode=True, content_analysis_mode=True, is_continuation=True, document_mode=True)
        else:
            # Стандартный режим анализа документов без вариантов
            logging.info("Standard content analysis mode without variants")
            await state.update_data(variant_mode=False, content_analysis_mode=True, is_continuation=True, document_mode=True)

    # Process and send the response to the user
    try:
        await process_and_send_ai_response(message, response, state=state)
    except Exception as e:
        # Log any errors that occur during response processing
        logging.error(f"Error processing AI response: {e}")
        import traceback
        logging.error(f"Traceback: {traceback.format_exc()}")
        # Try to send a basic error message to the user
        try:
            await message.answer("An error occurred while processing the content. Please try again later.")
        except Exception as msg_error:
            logging.error(f"Failed to send error message to user: {msg_error}")
    finally:
        # Clean up resources regardless of success or failure
        processing_users.discard(user_id)
        
        # Clear collected content for this user
        if user_id in collected_content:
            # Delete downloaded files
            for content_info in collected_content[user_id]:
                if 'file_path' in content_info:
                    try:
                        os.remove(content_info['file_path'])
                    except Exception as e:
                        logging.error(f"Error deleting file {content_info['file_path']}: {e}")
            
            # Clear the collected content
            collected_content[user_id] = []
