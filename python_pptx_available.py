"""
Module to check if python-pptx package is available
"""
import logging
import importlib.util
import sys

def is_python_pptx_available():
    """
    Check if python-pptx is installed and available.
    
    Returns:
        bool: True if python-pptx is available, False otherwise
    """
    try:
        # Check if package is installed
        if importlib.util.find_spec('pptx') is None:
            logging.warning("python-pptx package is not installed")
            return False
            
        # Try to import the package to ensure it works correctly
        import pptx
        
        # Try to create an empty presentation to validate the package works
        from pptx import Presentation
        temp_prs = Presentation()
        temp_slide = temp_prs.slides.add_slide(temp_prs.slide_layouts[0])
        
        # If we get here, the package is working
        logging.info("python-pptx package is available and working correctly")
        return True
        
    except ImportError as e:
        logging.error(f"Error importing python-pptx: {str(e)}")
        return False
    except Exception as e:
        logging.error(f"Unexpected error testing python-pptx: {str(e)}")
        return False
        
# If run directly, print the availability status
if __name__ == "__main__":
    available = is_python_pptx_available()
    print(f"python-pptx available: {available}")
