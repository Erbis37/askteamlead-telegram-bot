import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime
from typing import Dict, List, Optional
import io

async def create_project_presentation(project_info: Dict, trello_integration=None) -> io.BytesIO:
    """
    Create a professional project presentation for the team
    
    Args:
        project_info: Dictionary with project details
        trello_integration: Optional TrelloIntegration instance for real data
        
    Returns:
        BytesIO object containing the PPTX file
    """
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5]
    
    # Color scheme
    primary_color = RGBColor(44, 62, 107)  # Dark blue
    accent_color = RGBColor(255, 192, 0)  # Gold
    
    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = project_info.get('title', 'Project Overview')
    subtitle.text = f"Team Lead AI Presentation\n{datetime.now().strftime('%B %d, %Y')}"
    
    # Style title
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = primary_color
            run.font.size = Pt(44)
            run.font.bold = True
    
    # 2. Agenda Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Agenda"
    
    bullet_slide = slide.placeholders[1]
    bullet_slide.text = ("• Project Status Overview\n"
                        "• Current Sprint Progress\n"
                        "• Task Distribution\n"
                        "• Timeline & Milestones\n"
                        "• Risks & Dependencies\n"
                        "• Resource Allocation\n"
                        "• Next Steps")
    
    # 3. Project Status Overview
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Project Status Overview"
    
    if trello_integration:
        try:
            status = await trello_integration.get_project_status()
            content = (f"• Total Tasks: {status['total_tasks']}\n"
                      f"• Tasks in Progress: {status['by_column'].get('Doing', 0)}\n"
                      f"• Tasks in Review: {status['by_column'].get('Code Review', 0)}\n"
                      f"• Completed Tasks: {status['by_column'].get('Done', 0)}\n"
                      f"• Overdue Tasks: {len(status['overdue_tasks'])}\n"
                      f"• Tasks Without Deadlines: {len(status['no_deadline_tasks'])}")
        except:
            content = project_info.get('status_overview', 'Project status data unavailable')
    else:
        content = project_info.get('status_overview', 'Project status data unavailable')
    
    bullet_slide = slide.placeholders[1]
    bullet_slide.text = content
    
    # 4. Task Distribution by Team Member
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Task Distribution"
    
    if trello_integration and 'by_assignee' in locals() and status:
        content_lines = []
        for assignee, count in sorted(status['by_assignee'].items(), key=lambda x: x[1], reverse=True):
            content_lines.append(f"• {assignee}: {count} tasks")
        content = "\n".join(content_lines[:7])  # Show top 7
    else:
        content = project_info.get('task_distribution', '• Team member distribution data unavailable')
    
    bullet_slide = slide.placeholders[1]
    bullet_slide.text = content
    
    # 5. Sprint Progress
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Current Sprint Progress"
    
    sprint_content = project_info.get('sprint_progress', 
                                     "• Sprint Goal: Feature Development\n"
                                     "• Progress: 65% Complete\n"
                                     "• Days Remaining: 5\n"
                                     "• Velocity: On Track")
    
    bullet_slide = slide.placeholders[1]
    bullet_slide.text = sprint_content
    
    # 6. Key Risks & Mitigation
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Risks & Mitigation Strategies"
    
    risks = project_info.get('risks', [
        "• Technical Debt: Allocate 20% time for refactoring",
        "• Resource Availability: Cross-training team members",
        "• Scope Creep: Weekly stakeholder reviews",
        "• Integration Issues: Early testing with external APIs"
    ])
    
    bullet_slide = slide.placeholders[1]
    bullet_slide.text = "\n".join(risks) if isinstance(risks, list) else risks
    
    # 7. Next Steps
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Next Steps & Action Items"
    
    next_steps = project_info.get('next_steps', [
        "• Complete code reviews for pending PRs",
        "• Deploy feature X to staging environment",
        "• Schedule user testing sessions",
        "• Update documentation",
        "• Plan next sprint"
    ])
    
    bullet_slide = slide.placeholders[1]
    bullet_slide.text = "\n".join(next_steps) if isinstance(next_steps, list) else next_steps
    
    # 8. Thank You Slide
    slide = prs.slides.add_slide(title_only_layout)
    title = slide.shapes.title
    title.text = "Questions?"
    
    # Center the title
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add contact info
    txBox = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.text = "Generated by AI Team Lead Bot\nPowered by Claude AI"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output

async def create_architecture_presentation(architecture_info: Dict) -> io.BytesIO:
    """
    Create a technical architecture presentation
    
    Args:
        architecture_info: Dictionary with architecture details
        
    Returns:
        BytesIO object containing the PPTX file
    """
    prs = Presentation()
    
    # Similar structure but focused on technical architecture
    # Implementation would follow similar pattern as above
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output