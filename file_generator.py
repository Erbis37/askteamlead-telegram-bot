import logging
import os
import tempfile
import json
import random
import requests
import io
import re
import traceback
import time
import asyncio
from typing import Dict, List, Union, Optional, Any

# Required libraries for working with documents
# When deploying, make sure these libraries are installed
# pip install python-docx openpyxl pptx pandas reportlab icalendar pytz

from claude_utils import active_sessions, send_message
import anthropic
import logging
from typing import Dict, List, Union, Optional, Tuple


async def reformulate_file_request(original_text: str, file_type: str, user_id: int) -> str:
    """
    Reformulates an informal user request into a formal file creation request
    of a specific type, suitable for processing by the file generator.
    
    Args:
        original_text: Original user text
        file_type: Type of file to create (excel, pptx, doc, pdf, calendar)
        user_id: User ID for accessing session history
        
    Returns:
        Reformulated request for the file generator
    """
    # Mapping file types to corresponding request keywords in Russian
    file_type_to_request = {
        'excel': 'create excel table',
        'pptx': 'create powerpoint presentation',
        'doc': 'create word document',
        'pdf': 'create pdf document',
        'calendar': 'add to calendar'
    }
    
    # Get formal request for the corresponding file type
    formal_request = file_type_to_request.get(file_type, f"create {file_type} file")
    
    # Extract the intent and details from the user's message using Claude
    try:
        from claude_utils import CLAUDE_API_KEY, CLAUDE_MODEL
        
        # Create a system prompt for extracting structured information
        system_prompt = (
            f"Пользователь хочет создать файл типа {file_type}, но сделал неформальный запрос: '{original_text}'. "
            f"Мне нужно структурировать этот запрос для системы генерации файлов. "
            f"Начни структурированный запрос с фразы '{formal_request}', затем добавь все важные детали. "
            f"Например, если пользователь написал 'Добавь встречу с начальником в пятницу в мой календарь', ответ должен быть 'add to calendar встреча с начальником в пятницу'. "
            f"Это необходимо для того, чтобы система могла правильно извлечь все необходимые параметры."
        )
        
        # Format messages for the request - system prompt is a separate parameter
        messages = [
            {"role": "user", "content": original_text}
        ]
        
        # Send request to Claude
        client = anthropic.Anthropic(api_key=CLAUDE_API_KEY)
        response = client.messages.create(
            model=CLAUDE_MODEL,  # Use the same model as other requests
            max_tokens=200,
            system=system_prompt,  # System prompt as a separate parameter
            messages=messages,
            temperature=0.1
        )
        
        # Get the reformulated request
        reformulated_request = response.content[0].text.strip()
        logging.info(f"Reformulated request: '{original_text}' -> '{reformulated_request}'")
        return reformulated_request
        
    except Exception as e:
        logging.error(f"Error reformulating request: {str(e)}")
        # In case of error, just return the formal request with the original text
        return f"{formal_request} {original_text}"

async def detect_file_request(text: str) -> Optional[str]:
    """
    Determines if the text contains a request to create a file.
    
    Returns:
        File type (docx, xlsx, pdf, csv, pptx) or None if the request is not for file creation
    """
    if not text:
        return None
        
    text = text.lower()
    
    # Excel detection
    if any(keyword in text for keyword in ["сделай таблицу", "создай таблицу", "создай excel", "сделай excel", "excel файл", "таблицу excel"]):
        return "excel"
    
    # PowerPoint detection
    if any(keyword in text for keyword in ["сделай презентацию", "создай презентацию", "создай powerpoint", 
                                          "сделай powerpoint", "powerpoint файл", "презентацию powerpoint", 
                                          "создай ppt", "сделай ppt", "создай слайды", "сделай слайды", "pptx"]):
        return "pptx"
    
    # Calendar detection - with typo tolerance
    calendar_keywords = [
        "добавь в календарь", "доабвь в календарь", "добавить в календарь", "закинь в календарь",
        "добавь в келендарь", "добавь в колендарь", "доабвить в календарь", "в календарь добавь",
        "создай событие в календаре", "сделай календарную запись", "создай календарную запись",
        "создай ics", "сделай ics", "calendar", "ics файл", "напоминание в календаре", "создай встречу"
    ]
    
    # More flexible calendar detection using partial matching
    if any(keyword in text for keyword in calendar_keywords):
        return "calendar"
    
    # Even more flexible calendar detection - check for key terms appearing close together
    calendar_terms = ["календар", "келендар", "колендар"]
    add_terms = ["добав", "доабв", "доб", "закин", "событи", "встреч"]
    
    if any(cal_term in text for cal_term in calendar_terms):
        if any(add_term in text for add_term in add_terms):
            return "calendar"
    
    # Word document detection
    if any(keyword in text for keyword in ["создай документ", "сделай документ", "создай word", 
                                          "сделай word", "word файл", "документ word", "создай docx", "сделай docx"]):
        return "doc"
    
    # PDF detection
    if any(keyword in text for keyword in ["создай pdf", "сделай pdf", "pdf файл"]):
        return "pdf"
    
    # Add a fuzzy intent check using Claude for unclear requests
    # This is resource-intensive, so only use when needed for ambiguous requests
    if ("календар" in text or "встреч" in text or "событи" in text) and len(text) < 50:
        return "calendar"
        
    return None

async def get_image_from_url(url: str) -> Optional[bytes]:
    """
    Downloads an image from a URL
    
    Args:
        url: Image URL
        
    Returns:
        Image bytes or None if failed
    """
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return response.content
    except Exception as e:
        logging.error(f"Exception when fetching image: {str(e)}")
        return None

async def generate_file_content(user_message: str, conversation_history: List[Dict], file_type: str, user_id: int = None) -> Dict:
    """
    Generates file content based on user request and conversation history.
    
    Args:
        user_message: User's request message
        conversation_history: Previous conversation messages
        file_type: Type of file to generate
        user_id: User's ID for the session
        
    Returns:
        Dictionary with generated content
    """
    file_type_names = {
        'excel': 'Excel',
        'pptx': 'PowerPoint',
        'doc': 'Word',
        'pdf': 'PDF',
        'calendar': 'Calendar'
    }
    
    # Создаем системный промпт для Claude
    system_prompt = f"""
    Ты — ИИ-ассистент, который помогает создавать файлы формата {file_type.upper()} на основе запросов пользователя.
    Сгенерируй Python-код для создания файла {file_type.upper()}, который соответствует запросу пользователя.
    
    ИНСТРУКЦИИ:
    1. Создай полный, хорошо отформатированный файл {file_type.upper()}
    2. Используй соответствующие библиотеки (python-docx, openpyxl, python-pptx, reportlab и т.д.)
    3. Файл должен соответствовать языку, используемому в запросе пользователя
    4. Включи релевантное содержимое, форматирование и структуру
    5. Файл должен быть сразу готов к использованию
    """
    
    # Modify the standard prompt with full instructions for Claude
    modified_prompt = f"""
    Мне нужно создать файл {file_type.upper()} на основе предоставленного контента. Я выполню Python-код, который ты предоставишь, чтобы сгенерировать этот файл.

    КРИТИЧЕСКИ ВАЖНЫЕ ИНСТРУКЦИИ:
    1. Сгенерируй файл с ЗАМЕТНЫМ СОДЕРЖИМЫМ — не создавай пустые файлы или заглушки
    2. Включи полезный текст, данные и форматирование, подходящие для файла {file_type.upper()}
    3. Содержимое должно напрямую отражать наш разговор/предоставленный контент
    4. Сделай файл профессионально структурированным с правильным форматированием
    5. Убедись, что код включает ВСЕ необходимые импорты и зависимости
    6. Мысленно протестируй код, чтобы убедиться, что он работает без ошибок
    7. Добавь реальное, осмысленное содержимое даже для слайдов, ячеек и абзацев
    
    ОБЯЗАТЕЛЬНО ПРОВЕРЬ СВОЮ РАБОТУ:
    - Убедись, что твой код создает РЕАЛЬНОЕ СОДЕРЖИМОЕ, а не только структуру файла
    - Убедись, что весь текст/контент попадает в сам файл, а не только существует в коде
    - Дважды проверь, что твой подход добавляет значимое содержимое в файл
    - Убедись, что ты не создаешь пустой/минималистичный файл
    
    ЭТО КРАЙНЕ ВАЖНО: Мне нужен полный, насыщенный контентом файл, который можно сразу использовать.
        
    # КРИТИЧЕСКИЕ ПРОВЕРКИ:
    # 1. Всегда проверяй импорты в начале скрипта
    # 2. Добавляй правильную обработку исключений по всему коду
    # 3. Для PowerPoint: ВСЕГДА оборачивай загрузку изображений в try/except блоки
    # 4. Сохраняй файлы с явными путями, например '/tmp/telegram_bot_files/output.pptx'
    # 5. Выводи сообщения об успешном завершении для подтверждения создания файла

    Используй соответствующие библиотеки:
    - Для Excel: pandas и openpyxl
    - Для Word: python-docx
    - Для PDF: reportlab
    - Для PowerPoint: python-pptx
    - Для календаря: icalendar
    
    # ВАЖНО: При создании презентаций PowerPoint следуй этим рекомендациям:
    # 1. Всегда сначала проверяй импорты
    # 2. Используй правильную обработку ошибок
    # 3. Сохраняй файлы в /tmp/telegram_bot_files/
    # 4. Выводи сообщения об успехе/ошибке
    
    # Пример структуры для скрипта создания PowerPoint:
    """
    # Standard imports
    import os
    import sys
    import logging
    from datetime import datetime
    
    # Configure logging
    logging.basicConfig(level=logging.INFO)
    logger = logging.getLogger(__name__)
    
    # Verify required packages
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        logger.info("Successfully imported required packages")
    except ImportError as e:
        logger.error(f"Missing required package: {e}")
        logger.error("Please install with: pip install python-pptx")
        sys.exit(1)
    
    # Main function
    def create_presentation():
        try:
            prs = Presentation()
            
            # Add slides and content here
            # Example:
            # slide = prs.slides.add_slide(prs.slide_layouts[0])
            # title = slide.shapes.title
            # title.text = "Your Title Here"
            
            # Save the presentation
            output_path = '/tmp/telegram_bot_files/presentation.pptx'
            prs.save(output_path)
            print(f"SUCCESS: Presentation saved to {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            return False
    
    if __name__ == "__main__":
        if not create_presentation():
            sys.exit(1)
    """
    ```
    
    Important guidelines for the generated code:
    1. ALWAYS start with import verification
    2. Use proper error handling with try/except blocks
    3. Include logging for debugging
    4. Save files to /tmp/telegram_bot_files/ directory
    5. Print success/error messages to stdout/stderr
    6. Return proper exit codes (0 for success, 1 for error)
    
    {
        '''
    При создании презентации учти следующие требования:
    1. Презентация должна иметь титульный слайд
    2. Структурируй информацию на слайдах (не более 7-10 строк текста на слайд)
    3. Добавь подходящие заголовки слайдов
    4. Используй ясный, читаемый формат
    5. ОБЯЗАТЕЛЬНО добавь надежную обработку ошибок при загрузке изображений!
    6. Используй изображения Unsplash в презентации, но с надежной обработкой ошибок
    7. Всегда оборачивай загрузку изображений в try-except блок
    8. Если изображение не удается загрузить, продолжай создание презентации без него
    9. Подбирай релевантные ключевые слова на английском языке для поиска изображений
    10. Пример надежной работы с изображениями:
       ```python
       # Надежное добавление изображения в слайд с обработкой ошибок
       def add_image_from_url(slide, img_url, left, top, width, height):
           try:
               # Установи таймаут для запроса - не более 3 секунд
               from urllib.request import urlopen
               import socket
               from io import BytesIO
               import time
               import random
               
               # Добавляем случайный параметр к URL для предотвращения кеширования
               random_param = int(time.time() * 1000) + random.randint(1, 1000)
               if "?" in img_url:
                   img_url += f"&_r={random_param}"
               else:
                   img_url += f"?_r={random_param}"
               
               # Устанавливаем небольшой таймаут для запроса
               socket.setdefaulttimeout(3)
               
               # Пробуем загрузить изображение
               image_stream = urlopen(img_url)
               image_data = BytesIO(image_stream.read())
               slide.shapes.add_picture(image_data, left, top, width, height)
               print(f"Successfully added image from {img_url}")
               return True
           except Exception as e:
               # В случае любой ошибки, добавляем текстовый блок вместо изображения
               print(f"Error loading image {img_url}: {str(e)}")
               # Добавляем текстовый блок вместо изображения
               text_box = slide.shapes.add_textbox(left, top, width, height)
               text_frame = text_box.text_frame
               p = text_frame.add_paragraph()
               p.text = "[Изображение недоступно]"
               return False
       ```
    11. КРИТИЧЕСКИ ВАЖНО! В основном коде презентации ВСЕГДА используй функцию add_image_from_url с обработкой ошибок.
    ''' if file_type == 'pptx' else ''
    }
    
    {
        '''
    При создании календарного файла учти следующие требования:
    1. Используй библиотеку icalendar для создания файла .ics
    2. Извлеки из запроса всю необходимую информацию: название события, дату, время, место и описание
    3. Если в запросе не указана какая-то информация, сделай предположения исходя из контекста
    4. Обязательно укажи корректную timezone (обычно Europe/Moscow)
    ''' if file_type == 'calendar' else ''
    }
    
    Весь контент должен основываться на нашем разговоре.
    Дай мне полный исполняемый Python-код, который я смогу запустить.
    """
    
    # Prepare messages with conversation context and prompt
    from claude_utils import get_claude_completion
    final_messages = []
    
    # Format conversation history if available
    if conversation_history and isinstance(conversation_history, list) and len(conversation_history) > 0:
        # Make a copy of conversation history to avoid modifying the original
        history_messages = conversation_history.copy()
        
        # Log information about the history we're using
        message_count = len(history_messages)
        user_msgs = len([m for m in history_messages if m.get('role') == 'user'])
        assistant_msgs = len([m for m in history_messages if m.get('role') == 'assistant'])
        logging.info(f"Using conversation history for file generation: {message_count} total messages ({user_msgs} user, {assistant_msgs} assistant)")
        
        # Include conversation history in the request
        final_messages.extend(history_messages)
        
        # Add the current request with file generation prompt
        final_messages.append({"role": "user", "content": f"{user_message}\n\n{modified_prompt}"})
    else:
        # No history available, just use the single message with prompt
        logging.info("No conversation history available for file generation")
        final_messages.append({"role": "user", "content": f"{user_message}\n\n{modified_prompt}"})
    
    # Make request to Claude API without saving to history
    logging.info(f"Sending file generation request with {len(final_messages)} messages to Claude API")
    response = await get_claude_completion(
        messages=final_messages,
        user_id=user_id,
        max_tokens=4000,
        temperature=0.7,
        save_in_history=False  # Don't pollute chat history with technical prompts
    )
    
    # Split response into lines for processing
    file_lines = response.split("\n")
    
    # First, get a human-readable filename from Claude
    filename_prompt = f"""Based on the following content, generate a meaningful, human-readable filename in English. 
    The filename should be descriptive of the content, use only lowercase letters, numbers, and underscores.
    Include today's date (YYYY-MM-DD) and a short random ID at the end.
    Do NOT include any extension. Keep it under 30 characters. Return ONLY the filename, nothing else.
    
    CONTENT TO NAME: {user_message}
    """
    
    # Get a filename suggestion from Claude without affecting conversation history
    # But still use conversation context to generate a more relevant filename
    from claude_utils import get_claude_completion
    
    # Prepare filename generation messages with context
    filename_messages = []
    
    # Use limited conversation history (last 2-3 messages) for context
    # This is enough for filename generation without overwhelming the API
    if conversation_history and isinstance(conversation_history, list) and len(conversation_history) > 0:
        # Get the last few messages for context (max 3)
        relevant_history = conversation_history[-3:] if len(conversation_history) > 3 else conversation_history
        filename_messages.extend(relevant_history)
        logging.info(f"Using {len(relevant_history)} recent messages for filename context")
    
    # Add the filename generation prompt
    filename_messages.append({"role": "user", "content": filename_prompt})
    
    # Get filename suggestion
    filename_response = await get_claude_completion(
        messages=filename_messages,
        user_id=user_id,
        max_tokens=1000,
        temperature=0.7,
        save_in_history=False
    )
    
    # Clean up the suggested filename
    suggested_filename = filename_response.strip().lower()
    suggested_filename = re.sub(r'[^a-z0-9_-]', '_', suggested_filename)  # Replace invalid chars
    suggested_filename = re.sub(r'_+', '_', suggested_filename)  # Remove repeated underscores
    
    # Make sure it contains date and ID, add if missing
    import datetime
    today_date = datetime.datetime.now().strftime("%Y-%m-%d")
    import random
    random_id = f"{random.randint(1000, 9999)}"
    
    # Check if the date is already in the filename
    if today_date not in suggested_filename:
        suggested_filename = f"{suggested_filename}_{today_date}"
    
    # Check if there's some random ID-like pattern; if not, add one
    if not re.search(r'\d{3,}', suggested_filename):
        suggested_filename = f"{suggested_filename}_{random_id}"
    
    # Limit length
    if len(suggested_filename) > 30:
        suggested_filename = suggested_filename[:30]
    
    # Add appropriate extension
    file_name = suggested_filename
    if file_type == "excel":
        file_name += ".xlsx"
    elif file_type == "doc":
        file_name += ".docx"
    elif file_type == "pdf":
        file_name += ".pdf"
    elif file_type == "pptx":
        file_name += ".pptx"
    elif file_type == "calendar":
        file_name += ".ics"
    
    # If no filename found in response, generate one
    if not file_name:
        file_name = f"generated_{file_type}_{int(random.random() * 10000)}"
        # Add extension
        if file_type == "excel":
            file_name += ".xlsx"
        elif file_type == "doc":
            file_name += ".docx"
        elif file_type == "pdf":
            file_name += ".pdf"
        elif file_type == "pptx":
            file_name += ".pptx"
        elif file_type == "calendar":
            file_name += ".ics"
    
    # Extract Python code from response
    code_start = None
    code_end = None
    for i, line in enumerate(file_lines):
        if "```python" in line or "```" in line and code_start is None:
            code_start = i + 1
        elif "```" in line and code_start is not None:
            code_end = i
            break
    
    # If code block found, extract it
    python_code = ""
    if code_start is not None and code_end is not None:
        python_code = "\n".join(file_lines[code_start:code_end])
    else:
        # Fallback: try to extract all lines that look like Python code
        # This is a simple heuristic and may not work well
        potential_code_lines = []
        for line in file_lines:
            if "import " in line or "=" in line or "def " in line or "class " in line or line.strip().startswith("#"):
                potential_code_lines.append(line)
        if potential_code_lines:
            python_code = "\n".join(potential_code_lines)
    
    return {
        "file_name": file_name,
        "file_type": file_type,
        "code": python_code,
        "content": response,
        "user_request": user_message
    }

async def create_file(file_info: Dict) -> Optional[str]:
    """
    Creates a file using generated Python code and ensures files are not empty
    
    Args:
        file_info: Dictionary with file information
        
    Returns:
        Path to created file or None if creation failed
    """
    try:
        # Extract file parameters
        file_type = file_info.get('file_type', 'txt').lower()
        file_name = file_info.get('file_name', '')
        python_code = file_info.get('code', '')
        
        if not file_name:
            # Generate a filename if not provided
            import time
            import hashlib
            # Create a unique hash based on content
            content_hash = hashlib.md5(python_code.encode()).hexdigest()[:8]
            timestamp = int(time.time())
            file_name = f"generated_{file_type}_{timestamp}_{content_hash}.{file_type}"
        
        # Create the file path in a temporary directory
        file_dir = "/tmp/telegram_bot_files"
        os.makedirs(file_dir, exist_ok=True)
        
        if not file_name.endswith(f'.{file_type}') and not '.' in file_name:
            file_name = f"{file_name}.{file_type}"
            
        file_path = f"{file_dir}/{file_name}"
        
        # Calendar files are handled separately by the calendar generator module
        if file_type == 'calendar':
            raise ValueError("Calendar files should be handled by calendar_generator module")
        
        # Check if we have code to work with
        if not python_code or len(python_code.strip()) < 5:
            logging.error("No code provided for file creation")
            raise ValueError(f"Cannot create {file_type} file without code")
        
        # Clean up the Python code before saving
        import textwrap
        
        # Remove common leading whitespace from all lines
        cleaned_code = textwrap.dedent(python_code)
        
        # Remove any leading/trailing whitespace and ensure it ends with a newline
        cleaned_code = cleaned_code.strip() + '\n'
        
        # Save the cleaned Python code to a temporary file
        script_path = f"/tmp/telegram_bot_files/temp_script_{hash(cleaned_code)}.py"
        with open(script_path, 'w', encoding='utf-8') as f:
            f.write(cleaned_code)
            
        # Log the first few lines of the script for debugging
        logging.info(f"Script saved to {script_path}. First 10 lines:")
        with open(script_path, 'r', encoding='utf-8') as f:
            for i, line in enumerate(f.readlines()[:10]):
                logging.info(f"{i+1}: {line.rstrip()}")
        
        # Run the code in a separate process with enhanced error handling
        import subprocess
        logging.info(f"Executing Python script: {script_path} to create {file_type} file")
        
        # Print first few lines of the script for debugging
        script_preview = '\n'.join(python_code.split('\n')[:20]) + '\n...'
        logging.info(f"Script preview:\n{script_preview}")
        
        try:
            result = subprocess.run(['python', script_path], 
                                capture_output=True, 
                                text=True,
                                timeout=90)  # Increased timeout to 90 seconds
            
            # Log subprocess output
            if result.stdout:
                logging.info(f"Script output: {result.stdout[:500]}" + (
                    "..." if len(result.stdout) > 500 else ""))
            
            if result.returncode != 0:
                logging.error(f"Error creating file: {result.stderr}")
                
                # Additional diagnostic for PPTX files
                if file_type == 'pptx':
                    logging.error("Enhanced error diagnostics for PPTX file generation:")
                    # Check if specific error patterns are in the output
                    if "urllib.error.HTTPError" in result.stderr or "Service Unavailable" in result.stderr:
                        logging.error("Unsplash API error detected - service may be unavailable")
                    elif "ImportError" in result.stderr:
                        logging.error("Missing dependency. Please ensure python-pptx is installed")
                
                raise ValueError(f"Failed to create {file_type} file: {result.stderr[:200]}")
        except subprocess.TimeoutExpired:
            logging.error(f"Timeout executing script for {file_type} file generation")
            raise ValueError(f"Timeout while creating {file_type} file - script took too long to execute")
        
        # Extract file type from the file name
        file_type = file_name.split('.')[-1] if '.' in file_name else ''
        
        # Ensure the file has been created with enhanced file search
        if not os.path.exists(file_path):
            # Try to find the file by extension with a more thorough search
            import glob
            import time
            
            # Wait a brief moment for filesystem operations to complete
            time.sleep(0.5)
            
            # Search for recent files with the correct extension
            created_files = glob.glob(f"/tmp/telegram_bot_files/*.{file_type}")
            
            # Sort by creation time (most recent first)
            if created_files:
                created_files.sort(key=lambda x: os.path.getctime(x), reverse=True)
                found_file = created_files[0]
                logging.info(f"Found alternative file with matching extension: {found_file}")
                return found_file
                
            # For PowerPoint, try with both pptx and ppt extensions
            if file_type == 'pptx':
                alt_files = glob.glob("/tmp/telegram_bot_files/*.ppt")
                if alt_files:
                    alt_files.sort(key=lambda x: os.path.getctime(x), reverse=True)
                    found_file = alt_files[0]
                    logging.info(f"Found alternative PowerPoint file: {found_file}")
                    return found_file
            
            # List all files in the directory for diagnostics
            all_files = os.listdir("/tmp/telegram_bot_files/")
            recent_files = sorted(
                [(f, os.path.getctime(f"/tmp/telegram_bot_files/{f}")) for f in all_files],
                key=lambda x: x[1], reverse=True
            )[:5]
            
            logging.error(f"File {file_path} was not created by the script")
            logging.error(f"Most recent files in directory: {recent_files}")
            
            # Special handling for PowerPoint files
            if file_type == 'pptx':
                # Check if the package is available
                from python_pptx_available import is_python_pptx_available
                if not is_python_pptx_available():
                    logging.error("python-pptx package is not available - this is required for PowerPoint generation")
                    raise ValueError("Cannot create PowerPoint files: python-pptx package is not installed")
                
                # Try fallback method for PowerPoint creation
                try:
                    logging.info("Attempting fallback method for PowerPoint creation...")
                    # Create a simple PowerPoint with essential content
                    fallback_path = create_fallback_powerpoint(user_message=file_info.get('user_message', 'Presentation'), 
                                                           file_name=file_name)
                    if fallback_path and os.path.exists(fallback_path):
                        logging.info(f"Successfully created PowerPoint using fallback method: {fallback_path}")
                        return fallback_path
                    else:
                        logging.error("Fallback PowerPoint creation also failed")
                except Exception as fallback_error:
                    logging.error(f"Error in fallback PowerPoint creation: {str(fallback_error)}")
            
            # If we reach here, all methods failed
            raise ValueError(f"Failed to create {file_type} file")
        
            
        return file_path
            
    except Exception as e:
        logging.error(f"Exception when creating file: {str(e)}")
        raise e  # Re-raise to be handled by the caller

def create_fallback_powerpoint(user_message: str, file_name: str) -> Optional[str]:
    """
    Create a simple fallback PowerPoint presentation when the main generation method fails.
    This uses a direct approach with minimal dependencies to maximize reliability.
    
    Args:
        user_message: The user's original request message
        file_name: Suggested file name
        
    Returns:
        Path to the created file or None if creation failed
    """
    logging.info("Creating fallback PowerPoint presentation...")
    
    try:
        # Make sure the python-pptx package is available
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            logging.info("Successfully imported python-pptx")
        except ImportError as e:
            logging.error(f"Cannot create fallback PowerPoint - missing python-pptx: {str(e)}")
            return None
        
        # Create a basic presentation
        prs = Presentation()
        
        # Set output path
        output_dir = "/tmp/telegram_bot_files"
        os.makedirs(output_dir, exist_ok=True)
        
        # Ensure file has .pptx extension
        if not file_name.lower().endswith('.pptx'):
            file_name = f"{file_name.split('.')[0]}.pptx"
        
        output_path = os.path.join(output_dir, file_name)
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Презентация по запросу"
        
        # Add subtitle with user message (shortened if needed)
        subtitle = slide.placeholders[1]  # Subtitle placeholder
        short_message = user_message[:200] + '...' if len(user_message) > 200 else user_message
        subtitle.text = f"Тема: {short_message}"
        
        # Add content slide with user request details
        bullet_slide_layout = prs.slide_layouts[1]  # Bullet layout
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Set slide title
        title = slide.shapes.title
        title.text = "Содержание запроса"
        
        # Add content
        content = slide.placeholders[1]  # Content placeholder
        tf = content.text_frame
        
        # Add user message as bullet points (split by lines or sentences)
        import re
        # Split message into logical segments for bullet points
        segments = re.split(r'[.\n]', user_message)
        segments = [s.strip() for s in segments if s.strip()]
        
        # Add only the first 5-7 segments as bullets to avoid overcrowding
        for i, segment in enumerate(segments[:6]):
            if i == 0:
                # First bullet point is already there
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = segment
            p.level = 0
        
        # If we had to truncate, add a note
        if len(segments) > 6:
            p = tf.add_paragraph()
            p.text = "...и другие детали"    
            p.level = 0
        
        # Add conclusion slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title box
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(1.5)
        
        # Add title to the slide
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.text = "Спасибо за внимание"
        
        # Format the title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(44)
        run.font.bold = True
        
        # Add timestamp and metadata
        import datetime
        now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
        footer_box = slide.shapes.add_textbox(left, Inches(5), width, Inches(1))
        tf = footer_box.text_frame
        tf.text = f"Создано: {now}"
        
        # Save the presentation to the output file
        try:
            prs.save(output_path)
            logging.info(f"Successfully created fallback PowerPoint: {output_path}")
            
            # Verify the file was created
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                return output_path
            else:
                logging.error("Fallback PowerPoint file was not created properly")
                return None
                
        except Exception as save_error:
            logging.error(f"Error saving fallback PowerPoint: {str(save_error)}")
            
            # Try one more time with a different filename
            try:
                alt_path = os.path.join(output_dir, f"fallback_presentation_{int(datetime.datetime.now().timestamp())}.pptx")
                prs.save(alt_path)
                logging.info(f"Saved fallback PowerPoint with alternate path: {alt_path}")
                return alt_path
            except Exception as alt_save_error:
                logging.error(f"Error saving alternate fallback PowerPoint: {str(alt_save_error)}")
                return None
    
    except Exception as e:
        logging.error(f"Unexpected error in fallback PowerPoint creation: {str(e)}")
        return None

async def generate_file(user_message: str, file_type: str, user_id: int = None, conversation_history: List[Dict] = None) -> Optional[Union[str, List[str]]]:
    """
    Generate a file based on a user message and conversation history.
    
    Args:
        user_message: User's message with the file generation request
        file_type: Type of file to generate
        user_id: User's ID
        conversation_history: Previous conversation history
    
    Returns:
        Path to the generated file, list of paths for calendar files, or None if generation failed
    """
    try:
        logging.info(f"Starting file generation for type: {file_type}")
        
        # Special case for calendar files
        if file_type == 'calendar':
            from calendar_generator import generate_calendar_files
            event_files = await generate_calendar_files(user_message, conversation_history, user_id)
            if event_files and len(event_files) > 0:
                # Create base directory for files
                os.makedirs("/tmp/telegram_bot_files", exist_ok=True)
                result_paths = []
                
                # Create each platform file
                for platform_info in event_files:
                    # Extract platform-specific info
                    platform = platform_info.get('platform', 'generic')
                    filename = platform_info.get('filename')
                    
                    # Получаем содержимое файла, которое уже было создано в generate_calendar_files
                    file_content = platform_info.get('content')
                    
                    # Подробное логирование для диагностики
                    if file_content:
                        logging.info(f"File content available for {platform}, length: {len(file_content)} bytes")
                    else:
                        logging.warning(f"No file content found for {platform}, attempting to create from event_details")
                    
                    # Если content отсутствует, но есть event_details, создаем файл на месте
                    if not file_content and 'event_details' in platform_info:
                        from calendar_generator import create_ics_file
                        event_details = platform_info.get('event_details', {})
                        
                        # Проверка наличия обязательных полей в event_details
                        if not event_details.get('summary'):
                            event_details['summary'] = "Событие" # Базовое название, если нет в ответе
                            logging.warning(f"Missing summary in event_details for {platform}, using default")
                        
                        logging.info(f"Creating ICS file from event_details: {json.dumps(event_details, ensure_ascii=False)[:200]}...")
                        
                        # Создаем содержимое календарного файла
                        file_content = create_ics_file(event_details)
                        logging.info(f"Generated file content, length: {len(file_content) if file_content else 0} bytes")
                    
                    # Write to disk only if we have real content (prevent empty files)
                    if file_content and len(file_content) > 10 and filename:
                        file_path = f"/tmp/telegram_bot_files/{filename}"
                        
                        # Сохраняем в файл и проверяем результат
                        with open(file_path, 'wb') as f:
                            f.write(file_content)
                        
                        # Проверяем, что файл действительно создан и имеет содержимое
                        if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                            result_paths.append(file_path)
                            logging.info(f"Successfully created calendar file for {platform}: {file_path} ({os.path.getsize(file_path)} bytes)")
                        else:
                            logging.error(f"Created file is empty or not found: {file_path}")
                    else:
                        logging.error(f"Skipping file creation for {platform}: insufficient content or missing filename")
                        if not filename:
                            logging.error("Missing filename")
                        if not file_content or len(file_content) <= 10:
                            logging.error(f"File content too small: {len(file_content) if file_content else 0} bytes")
                
                return result_paths
            else:
                logging.error("Calendar file generation failed")
                raise Exception("Failed to generate calendar event details")
        
        # Generate file content using AI - similar to how calendar generator works
        file_info = await generate_file_content(user_message, conversation_history or [], file_type, user_id)
        
        if not file_info:
            logging.error("Failed to generate file content")
            raise Exception(f"Failed to generate content for {file_type} file")
            
        # Create file using generated code
        file_path = await create_file(file_info)
        
        if not file_path or not os.path.exists(file_path):
            logging.error("File creation failed")
            raise Exception(f"Failed to create {file_type} file")
            
        return file_path
    except Exception as e:
        error_message = str(e)
        logging.error(f"File generation error: {error_message}")
        # Raise the exception to be handled by the caller
        raise e

async def handle_file_generation_error(file_type: str, error_message: str) -> str:
    """
    Handles file generation errors by providing a user-friendly error message.
    
    Args:
        file_type: Type of file that failed to generate
        error_message: Detailed error message
        
    Returns:
        User-friendly error message
    """
    file_type_names = {
        'doc': 'документа Word',
        'docx': 'документа Word',
        'excel': 'таблицы Excel',
        'xlsx': 'таблицы Excel',
        'pptx': 'презентации PowerPoint',
        'pdf': 'PDF документа',
        'calendar': 'события календаря'
    }
    
    file_type_display = file_type_names.get(file_type, f'файла {file_type}')
    
    # Log the full error for debugging
    logging.error(f"File generation error ({file_type}): {error_message}")
    
    # Create a user-friendly error message
    user_message = (
        f"Не удалось создать файл {file_type_display}. "
        f"Пожалуйста, проверьте запрос и попробуйте еще раз с более подробной информацией."
    )
    
    return user_message

async def extract_event_info(content: str, conversation_context: List[Dict] = None, user_id: int = None) -> Dict:
    try:
        from calendar_generator import generate_calendar_files
        
        # Get the event details from the first platform (they're all the same content-wise)
        calendar_files = await generate_calendar_files(content, conversation_context)
        
        if calendar_files and len(calendar_files) > 0:
            # Extract the event details from the first platform
            first_file = calendar_files[0]
            if 'event_details' in first_file:
                return first_file['event_details']
        
        # If that fails, return minimal info
        return {
            "summary": "Event",
            "start_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
            "end_time": (datetime.datetime.now() + datetime.timedelta(hours=1)).strftime("%Y-%m-%d %H:%M"),
            "location": "",
            "description": content
        }
        
    except Exception as e:
        logging.error(f"Error extracting event info: {str(e)}")
        # Return minimal information on error
        return {
            "summary": "Event",
            "start_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
            "end_time": (datetime.datetime.now() + datetime.timedelta(hours=1)).strftime("%Y-%m-%d %H:%M"),
            "location": "",
            "description": content
        }

async def create_calendar_file(file_path: str, content: str, code: str, user_id: int = None) -> Optional[List[str]]:
    """
    Create multiple iCalendar files for different platforms (iOS, Android, Outlook)
    by delegating to the specialized calendar_generator module
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code (not used in this implementation)
        user_id: User ID for session management
    
    Returns:
        List of paths to the generated files or None if creation failed
    """
    try:
        # Get conversation context if available
        conversation_context = None
        if user_id and user_id in active_sessions:
            session = active_sessions[user_id]
            # Fix: Use 'messages' key instead of 'history' to match how it's stored in claude_utils.py
            conversation_context = session.get("messages", [])
            logging.info(f"Using existing conversation context for calendar generation, user_id={user_id}, message count: {len(conversation_context)}")
            # Log the first few messages to help with debugging
            if conversation_context:
                logging.info(f"First message in context: {str(conversation_context[0])[:100]}...")
        
        # Import the new calendar generator functionality
        from calendar_generator import generate_calendar_files
        
        # Generate calendar files using the new implementation
        # Pass the user_id explicitly to ensure proper context handling
        logging.info(f"Passing user_id={user_id} to calendar generator")
        result_files = await generate_calendar_files(content, conversation_context, user_id=user_id)
        
        if not result_files:
            logging.error("Failed to create calendar files: No files returned from generator")
            return None
        
        # We need to save the files to disk since the new implementation returns file content in memory
        saved_files = []
        
        # Create destination directory if it doesn't exist
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        # Save each file
        for file_info in result_files:
            filename = file_info.get('filename')
            content = file_info.get('content')
            
            if not filename or not content:
                continue
                
            # Create full path for the file
            full_path = os.path.join(os.path.dirname(file_path), filename)
            
            # Save the file
            with open(full_path, 'wb') as f:
                f.write(content)
                
            saved_files.append(full_path)
            logging.info(f"Saved calendar file: {full_path}")
        
        # Status messages are now handled directly in message_handlers.py
        # using working_msg.delete()
        if saved_files:
            return saved_files
        else:
            logging.error("No calendar files were successfully saved")
            return None
            
    except Exception as e:
        error_info = traceback.format_exc()
        logging.error(f"Failed to create calendar file: {str(e)}\nTraceback: {error_info}")
        return None

async def create_doc_file(file_path: str, content: str, code: str, user_id: int = None) -> Optional[str]:
    """
    Create a Word document (DOCX) based on user request and conversation context
    using AI to generate the content.
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code for creating the document
        user_id: User ID for session management
    
    Returns:
        Path to the generated file or None if creation failed
    """
    try:
        # Create destination directory if it doesn't exist
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        # Check if code contains valid Python code for document creation
        if code and 'Document(' in code:
            # Create a temporary Python file to execute the document creation code
            script_path = f"{os.path.dirname(file_path)}/doc_generator.py"
            
            # Write the code to a script file
            with open(script_path, 'w') as f:
                f.write(code)
                
            try:
                # Prepare environment variables for the script
                env = os.environ.copy()
                env['OUTPUT_PATH'] = file_path
                
                # Run the script to generate the document
                process = await asyncio.create_subprocess_exec(
                    'python', script_path,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    env=env
                )
                
                stdout, stderr = await process.communicate()
                
                # Check if the file was successfully created
                if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                    # Delete the temporary script
                    try:
                        os.remove(script_path)
                    except Exception as e:
                        logging.error(f"Error removing temporary script: {e}")
                        
                    return file_path
                else:
                    logging.error(f"Document generation failed or created empty file. Code execution output: {stdout.decode()}\nErrors: {stderr.decode()}")
            except Exception as script_error:
                logging.error(f"Error executing document script: {str(script_error)}")
                
        # If we reach here, code execution failed or didn't generate a valid file
        # Use fallback method
        return await create_doc_fallback(file_path, content, code)
        
    except Exception as e:
        logging.error(f"Failed to create Word document: {str(e)}\nTraceback: {traceback.format_exc()}")
        return None


async def create_doc_fallback(file_path: str, content: str, code: str) -> Optional[str]:
    """
    Create a Word document as fallback when the generated code fails.
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code (not used in fallback)
    
    Returns:
        Path to the generated file or None if creation failed
    """
    try:
        from docx import Document
        document = Document()
        
        # Add title extracted from content
        title = content.split('\n')[0] if content else "Generated Document"
        document.add_heading(title, 0)
        
        # Add content broken into paragraphs
        paragraphs = content.split('\n\n') if content else ["Sample document content"]
        for paragraph in paragraphs:
            if paragraph.strip():
                document.add_paragraph(paragraph)
        
        # Save the document
        document.save(file_path)
        return file_path
    except Exception as e:
        logging.error(f"Error creating Word fallback: {str(e)}")
        return None


async def create_excel_file(file_path: str, content: str, code: str, user_id: int = None) -> Optional[str]:
    """
    Create an Excel spreadsheet (XLSX) based on user request and conversation context
    using AI to generate the content.
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code for creating the Excel file
        user_id: User ID for session management
    
    Returns:
        Path to the generated file or None if creation failed
    """
    try:
        # Ensure the directory exists
        os.makedirs(os.path.dirname(os.path.abspath(file_path)), exist_ok=True)
        logging.info(f"Creating Excel file at: {os.path.abspath(file_path)}")
        
        # First, try to create the Excel file directly using pandas
        try:
            import pandas as pd
            import json
            
            # Try to parse the content as JSON first
            try:
                data = json.loads(content)
                if isinstance(data, dict):
                    # If it's a dictionary, create a DataFrame with one row
                    df = pd.DataFrame([data])
                elif isinstance(data, list):
                    # If it's a list, create a DataFrame with multiple rows
                    df = pd.DataFrame(data)
                else:
                    raise ValueError("Unsupported JSON format")
                
                # Ensure the file has .xlsx extension
                if not file_path.lower().endswith('.xlsx'):
                    file_path = f"{os.path.splitext(file_path)[0]}.xlsx"
                
                # Write to Excel file with explicit engine
                df.to_excel(file_path, index=False, engine='openpyxl')
                
                # Verify the file was created
                if os.path.exists(file_path):
                    file_size = os.path.getsize(file_path)
                    logging.info(f"Excel file created successfully. Size: {file_size} bytes")
                    if file_size > 0:
                        return file_path
                    else:
                        logging.error("Created Excel file is empty")
                else:
                    logging.error("Excel file was not created")
                
            except (json.JSONDecodeError, ValueError) as e:
                logging.info(f"Content is not JSON, trying tabular format: {str(e)}")
                # If JSON parsing fails, try tabular format
                return await create_excel_fallback(file_path, content, code)
            
        except ImportError as ie:
            logging.error(f"Required package not found: {str(ie)}")
            return None
            
        except Exception as direct_error:
            logging.error(f"Direct Excel creation failed: {str(direct_error)}\n{traceback.format_exc()}")
            return await create_excel_fallback(file_path, content, code)
            
    except Exception as e:
        logging.error(f"Unexpected error in create_excel_file: {str(e)}\n{traceback.format_exc()}")
        try:
            # One last attempt with minimal content
            import pandas as pd
            
            # Ensure the directory exists
            os.makedirs(os.path.dirname(os.path.abspath(file_path)), exist_ok=True)
            
            # Ensure .xlsx extension
            if not file_path.lower().endswith('.xlsx'):
                file_path = f"{os.path.splitext(file_path)[0]}.xlsx"
            
            # Create a simple dataframe
            df = pd.DataFrame({
                'Error': ['Failed to generate requested content'], 
                'Message': ['Here is your original content instead:'],
                'Content': [content[:1000] + ('...' if len(content) > 1000 else '')]
            })
            
            # Write with explicit engine
            df.to_excel(file_path, index=False, engine='openpyxl')
            
            if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                return file_path
            return None
            
        except Exception as final_error:
            logging.error(f"Final fallback failed: {str(final_error)}\n{traceback.format_exc()}")
            return None


async def create_excel_fallback(file_path: str, content: str, code: str) -> Optional[str]:
    """
    Create an Excel file as fallback when the generated code fails.
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code (not used in fallback)
    
    Returns:
        Path to the generated file or None if creation failed
    """
    try:
        # Ensure the directory exists
        os.makedirs(os.path.dirname(os.path.abspath(file_path)), exist_ok=True)
        
        # Ensure .xlsx extension
        if not file_path.lower().endswith('.xlsx'):
            file_path = f"{os.path.splitext(file_path)[0]}.xlsx"
        
        import pandas as pd
        import re
        
        # If content is empty, use a default message
        if not content.strip():
            content = "No content provided"
        
        # Process content to extract table data
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        # If no lines, create a simple Excel file with the content
        if not lines:
            df = pd.DataFrame({'Content': [content[:1000] + ('...' if len(content) > 1000 else '')]})
            df.to_excel(file_path, index=False, engine='openpyxl')
            return file_path if os.path.exists(file_path) and os.path.getsize(file_path) > 0 else None
        
        # Try to determine if there's a header row
        header_row = lines[0] if len(lines) > 1 else None
        data_rows = lines[1:] if len(lines) > 1 else lines
        
        # Detect separator (tab, comma, or pipe)
        if header_row:
            if '\t' in header_row:
                sep = '\t'
            elif ',' in header_row:
                sep = ','
            elif '|' in header_row:
                sep = '|'
            else:
                # Default to comma if no clear separator
                sep = ','
                
            # Clean header if using pipe separator (common in markdown tables)
            if sep == '|':
                header_row = header_row.strip('|').strip()
                data_rows = [row.strip('|').strip() for row in data_rows if row.strip()]
                # Skip separator line in markdown tables (like |---|---|)
                if data_rows and re.match(r'^[-:|\s]+$', data_rows[0].replace('|', '')):
                    data_rows = data_rows[1:]
            
            # Process headers
            headers = [h.strip() for h in header_row.split(sep)]
            
            # Process data rows
            data = []
            for row in data_rows:
                if row.strip():
                    data.append([cell.strip() for cell in row.split(sep)])
            
            # Create DataFrame
            df = pd.DataFrame(data)
            
            # Try to use headers if the number matches
            if len(headers) == len(df.columns):
                df.columns = headers
            
            # Ensure all values are strings to avoid serialization issues
            df = df.astype(str)
            
            # Write to Excel file with explicit engine
            df.to_excel(file_path, index=False, engine='openpyxl')
            
            if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                return file_path
            
        # If we get here, try a simple approach with the content
        try:
            # Try to split into columns by whitespace
            data = [line.split() for line in lines if line.strip()]
            if data:
                df = pd.DataFrame(data[1:], columns=data[0] if len(data) > 1 else [f'Column {i+1}' for i in range(len(data[0]))])
                df.to_excel(file_path, index=False, engine='openpyxl')
                if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                    return file_path
        except Exception as simple_error:
            logging.warning(f"Simple table creation failed: {str(simple_error)}")
        
        # If all else fails, create a simple Excel file with the content
        df = pd.DataFrame({
            'Content': [content[:1000] + ('...' if len(content) > 1000 else '')],
            'Message': ['Failed to parse content as table. Showing raw content instead.']
        })
        df.to_excel(file_path, index=False, engine='openpyxl')
        return file_path if os.path.exists(file_path) and os.path.getsize(file_path) > 0 else None
        
    except ImportError as ie:
        logging.error(f"Required package not found in fallback: {str(ie)}")
        return None
        
    except Exception as e:
        logging.error(f"Error in Excel fallback: {str(e)}\n{traceback.format_exc()}")
        return None


async def create_pdf_file(file_path: str, content: str, code: str, user_id: int = None) -> Optional[str]:
    """
    Create a PDF document based on user request and conversation context
    using AI to generate the content.
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code for creating the PDF
        user_id: User ID for session management
    
    Returns:
        Path to the generated file or None if creation failed
    """
    try:
        # Create destination directory if it doesn't exist
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        # Check if code contains valid Python code for PDF creation
        if code and 'reportlab' in code:
            # Validate Python code syntax first
            try:
                compile(code, '<string>', 'exec')
            except SyntaxError as e:
                logging.error(f"Invalid Python syntax in generated code: {str(e)}")
                # Fall back to the fallback method if syntax is invalid
                return await create_pdf_fallback(file_path, content, code)
            
            # Create a temporary Python file to execute the PDF creation code
            script_path = f"{os.path.dirname(file_path)}/pdf_generator_{int(time.time())}.py"
            
            try:
                # Write the code to a script file with error handling
                with open(script_path, 'w', encoding='utf-8') as f:
                    f.write(code)
                
                # Prepare environment variables for the script
                env = os.environ.copy()
                env['OUTPUT_PATH'] = file_path
                
                # Run the script to generate the PDF with a timeout
                process = await asyncio.create_subprocess_exec(
                    'python', script_path,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    env=env
                )
                
                # Set a timeout for the process (30 seconds)
                try:
                    stdout, stderr = await asyncio.wait_for(process.communicate(), timeout=30.0)
                except asyncio.TimeoutError:
                    process.terminate()
                    await asyncio.sleep(1)
                    if process.returncode is None:
                        process.kill()
                    raise RuntimeError("PDF generation timed out after 30 seconds")
                
                # Check if the file was successfully created
                if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                    return file_path
                else:
                    error_output = stderr.decode('utf-8', errors='replace')
                    logging.error(f"PDF generation failed or created empty file. Errors: {error_output}")
            except Exception as script_error:
                logging.error(f"Error executing PDF script: {str(script_error)}\n{traceback.format_exc()}")
            finally:
                # Always clean up the temporary script
                try:
                    if os.path.exists(script_path):
                        os.remove(script_path)
                except Exception as e:
                    logging.error(f"Error removing temporary script {script_path}: {e}")
                
        # If we reach here, code execution failed or didn't generate a valid file
        # Use fallback method
        return await create_pdf_fallback(file_path, content, code)
        
    except Exception as e:
        logging.error(f"Failed to create PDF file: {str(e)}\nTraceback: {traceback.format_exc()}")
        return None


async def create_pdf_fallback(file_path: str, content: str, code: str) -> Optional[str]:
    """
    Create a PDF document as fallback when the generated code fails.
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code (not used in fallback)
    
    Returns:
        Path to the generated file or None if creation failed
    """
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_LEFT, TA_CENTER
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.lib import fonts
        
        # Register standard fonts explicitly
        pdfmetrics.registerFontFamily(
            'Helvetica',
            normal='Helvetica',
            bold='Helvetica-Bold',
            italic='Helvetica-Oblique',
            boldItalic='Helvetica-BoldOblique'
        )
        
        # Import required font handling modules
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont, CIDFont
        from reportlab.pdfbase.pdfmetrics import registerFontFamily
        
        # Use a built-in font that supports Cyrillic
        try:
            # Try to register a font that supports Cyrillic
            try:
                # Try DejaVuSans if available (common on Linux)
                font_path = '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'
                pdfmetrics.registerFont(TTFont('DejaVuSans', font_path))
                default_font = 'DejaVuSans'
                logging.info("Using DejaVuSans font for Cyrillic support")
            except:
                # Fall back to Arial Unicode MS (common on Windows)
                try:
                    font_path = 'Arial Unicode MS'
                    pdfmetrics.registerFont(TTFont('ArialUnicode', font_path))
                    default_font = 'ArialUnicode'
                    logging.info("Using Arial Unicode MS font for Cyrillic support")
                except:
                    # As last resort, use the built-in Helvetica
                    default_font = 'Helvetica'
                    logging.warning("Using default Helvetica font, Cyrillic support may be limited")
        except Exception as e:
            logging.error(f"Error setting up fonts: {str(e)}")
            default_font = 'Helvetica'
        
        # Create document with proper margins
        doc = SimpleDocTemplate(
            file_path,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72,
            encoding='utf-8'
        )
        
        # Get default styles
        styles = getSampleStyleSheet()
        
        # Update existing styles for Russian text support
        normal_style = styles['Normal']
        normal_style.fontName = default_font
        normal_style.fontSize = 10
        normal_style.leading = 14
        normal_style.encoding = 'utf-8'
        
        title_style = styles['Title']
        title_style.fontName = f'{default_font}-Bold' if default_font != 'Helvetica' else 'Helvetica-Bold'
        title_style.fontSize = 18
        title_style.leading = 22
        title_style.alignment = TA_CENTER
        title_style.spaceAfter = 20
        title_style.encoding = 'utf-8'
        
        heading1_style = styles['Heading1']
        heading1_style.fontName = f'{default_font}-Bold' if default_font != 'Helvetica' else 'Helvetica-Bold'
        heading1_style.fontSize = 16
        heading1_style.leading = 20
        heading1_style.spaceAfter = 12
        heading1_style.encoding = 'utf-8'
        
        # Build the content
        story = []
        
        try:
            # Extract title from first line of content
            lines = content.split('\n')
            title = lines[0] if lines else "Сгенерированный PDF"
            
            # Add title with error handling
            try:
                title_para = Paragraph(title, title_style)
                story.append(title_para)
                story.append(Spacer(1, 12))
            except Exception as e:
                logging.warning(f"Error adding title to PDF: {str(e)}")
                # Try with simplified title if there's an encoding issue
                title = "Сгенерированный документ"
                story.append(Paragraph(title, styles['Title']))
                story.append(Spacer(1, 12))
            
            # Process content by paragraphs with error handling
            paragraphs = content.split('\n\n')
            for i, paragraph in enumerate(paragraphs):
                if not paragraph.strip():
                    continue
                    
                try:
                    # Clean up the paragraph text and ensure proper encoding
                    clean_para = paragraph.replace('\n', '<br/>')
                    # Ensure proper handling of special characters and emojis
                    clean_para = clean_para.encode('utf-8', errors='replace').decode('utf-8')
                    # Replace any remaining problematic characters
                    clean_para = ''.join(c if ord(c) < 0x10000 else '?' for c in clean_para)
                    # Add paragraph with proper style
                    p = Paragraph(clean_para, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 6))
                except Exception as para_error:
                    logging.warning(f"Error adding paragraph {i} to PDF: {str(para_error)}")
                    continue
            
            # Build the PDF
            doc.build(story)
            
            # Verify the file was created
            if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                return file_path
            return None
            
        except Exception as content_error:
            logging.error(f"Error building PDF content: {str(content_error)}")
            return None
            
    except Exception as e:
        logging.error(f"Error in create_pdf_fallback: {str(e)}\n{traceback.format_exc()}")
        return None


async def create_pptx_file(file_path: str, content: str, code: str, user_id: int = None) -> Optional[str]:
    """
    Create a PowerPoint presentation (PPTX) based on user request and conversation context
    using AI to generate the content with Unsplash images.
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code for creating the presentation
        user_id: User ID for session management
    
    Returns:
        Path to the generated file or None if creation failed
    """
    try:
        # Create destination directory if it doesn't exist
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        # Check if code contains valid Python code for PowerPoint creation
        if code and 'pptx' in code:
            # Clean the code to handle string literals properly
            try:
                # Fix common string literal issues
                code = code.replace('\n', '\\n').replace('\"', '\\"').replace("'", "\\'")
                
                # Create a temporary Python file to execute the presentation creation code
                script_path = f"{os.path.dirname(file_path)}/pptx_generator_{int(time.time())}.py"
                
                # Write the code to a script file with proper encoding
                with open(script_path, 'w', encoding='utf-8') as f:
                    # Add encoding header and error handling
                    f.write('#!/usr/bin/env python\n')
                    f.write('# -*- coding: utf-8 -*-\n\n')
                    f.write('import sys\n')
                    f.write('import os\n\n')
                    f.write('def main():\n')
                    f.write('    try:\n')
                    f.write('        output_path = os.environ.get(\'OUTPUT_PATH\')\n')
                    f.write('        # The actual code with proper string handling\n')
                    
                    # Split code into lines and indent them
                    for line in code.split('\n'):
                        f.write(f'        {line}\n')
                        
                    f.write('        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:\n')
                    f.write('            print(f"Successfully created presentation at {output_path}")\n')
                    f.write('            return 0\n')
                    f.write('        else:\n')
                    f.write('            print("Error: Failed to create presentation", file=sys.stderr)\n')
                    f.write('            return 1\n')
                    f.write('    except Exception as e:\n')
                    f.write('        print(f"Error in presentation generation: {str(e)}", file=sys.stderr)\n')
                    f.write('        return 1\n\n')
                    f.write('if __name__ == \'__main__\':\n')
                    f.write('    sys.exit(main())\n')
                
                # Make the script executable
                os.chmod(script_path, 0o755)
                
                # Prepare environment variables for the script
                env = os.environ.copy()
                env['PYTHONIOENCODING'] = 'utf-8'
                env['OUTPUT_PATH'] = file_path
                
                # Run the script to generate the presentation
                process = await asyncio.create_subprocess_exec(
                    'python3', script_path,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    env=env
                )
                
                # Set a timeout for the process (60 seconds)
                try:
                    stdout, stderr = await asyncio.wait_for(process.communicate(), timeout=60.0)
                    stdout = stdout.decode('utf-8', errors='replace')
                    stderr = stderr.decode('utf-8', errors='replace')
                    
                    # Check if the file was successfully created
                    if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                        logging.info(f"Successfully created presentation: {stdout}")
                        return file_path
                    else:
                        error_msg = f"PowerPoint generation failed. Output: {stdout}\nErrors: {stderr}"
                        logging.error(error_msg)
                        
                except asyncio.TimeoutError:
                    process.terminate()
                    await asyncio.sleep(1)
                    if process.returncode is None:
                        process.kill()
                    error_msg = "PowerPoint generation timed out after 60 seconds"
                    logging.error(error_msg)
                    
            except Exception as script_error:
                error_msg = f"Error in PowerPoint script generation: {str(script_error)}\n{traceback.format_exc()}"
                logging.error(error_msg)
            finally:
                # Clean up the temporary script
                try:
                    if os.path.exists(script_path):
                        os.remove(script_path)
                except Exception as e:
                    logging.error(f"Error removing temporary script: {e}")
                
        # If we reach here, code execution failed or didn't generate a valid file
        # Use fallback method
        return await create_pptx_fallback(file_path, content, code)
        
    except Exception as e:
        logging.error(f"Failed to create PowerPoint presentation: {str(e)}\nTraceback: {traceback.format_exc()}")

async def create_pptx_fallback(file_path: str, content: str, code: str) -> Optional[str]:
    """
    Create a PowerPoint presentation as fallback when the generated code fails.
    Uses a simple approach to ensure reliability.
    
    Args:
        file_path: Path where the file should be saved
        content: User's message content
        code: Generated Python code (not used in fallback)
    
    Returns:
        Path to the generated file or None if creation failed
    """
    try:
        # Try to import required modules
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError as e:
            logging.error(f"Required PowerPoint modules not available: {e}")
            return None
        
        try:
            # Create a new presentation
            prs = Presentation()
            
            # Add a title slide
            title_slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            # Set title and subtitle
            title.text = "Your Presentation"
            subtitle.text = f"Created on {datetime.now().strftime('%Y-%m-%d')}"
            
            # Split content into slides (one slide per paragraph)
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            
            # If no paragraphs found, use the whole content
            if not paragraphs:
                paragraphs = [content]
            
            # Create content slides (limit to 10 slides max)
            for i, para in enumerate(paragraphs[:10]):
                try:
                    # Use a content layout for the slide
                    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    
                    # Set slide title
                    title = content_slide.shapes.title
                    title.text = f"Slide {i+1}"
                    
                    # Add content
                    body_shape = content_slide.shapes.placeholders[1]
                    tf = body_shape.text_frame
                    
                    # Add the paragraph text
                    p = tf.add_paragraph()
                    p.text = para
                    p.level = 0
                    
                except Exception as slide_error:
                    logging.error(f"Error creating slide {i+1}: {str(slide_error)}")
                    continue
            
            # Add a closing slide
            try:
                end_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                left = Inches(1)
                top = Inches(2)
                width = prs.slide_width - Inches(2)
                height = Inches(2)
                
                textbox = end_slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                
                p = tf.add_paragraph()
                p.text = "Thank You!"
                p.font.size = Pt(44)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            except Exception as end_error:
                logging.error(f"Error creating end slide: {str(end_error)}")
            
            # Ensure the directory exists
            os.makedirs(os.path.dirname(file_path), exist_ok=True)
            
            # Save the presentation
            prs.save(file_path)
            
            # Verify the file was created and is not empty
            if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                logging.info(f"Successfully created fallback presentation at {file_path}")
                return file_path
            else:
                logging.error("Failed to create fallback presentation: File is empty or not created")
                return None
                
        except Exception as e:
            logging.error(f"Error in create_pptx_fallback: {str(e)}", exc_info=True)
            return None
            
    except Exception as e:
        logging.error(f"Unexpected error in create_pptx_fallback: {str(e)}", exc_info=True)
        return None

async def process_file_request(message_text: str, conversation_history: List[Dict]) -> Optional[Dict]:
    """
    Process a request to generate a file from the user message.
    
    Args:
        message_text: User message text
        conversation_history: History of the conversation
        
    Returns:
        Dictionary with file information or None if not a file request
    """
    file_type = await detect_file_request(message_text)
    
    if not file_type:
        return None
        
    return {
        "file_type": file_type,
        "message": message_text,
        "conversation": conversation_history
    }
